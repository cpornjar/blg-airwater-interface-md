"""
generate_pptx.py — ACADEMIC_NAVY themed PPTX for Keynote import
================================================================
Generates slides/presentation.pptx with navy/gold academic theme.
Open in Keynote on macOS → File → Save As → .key

Usage:  python3 slides/generate_pptx.py
Output: slides/presentation.pptx
"""

from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from pptx.oxml.ns import qn
import copy

# ── Paths ──────────────────────────────────────────────────────────────────────
ROOT   = Path(__file__).resolve().parent
FIGS   = ROOT / "figures"
OUT    = ROOT / "presentation.pptx"

# ── Colors ─────────────────────────────────────────────────────────────────────
NAVY   = RGBColor(0x1A, 0x27, 0x44)
GOLD   = RGBColor(0xC9, 0xA8, 0x4C)
BLUE   = RGBColor(0x44, 0x72, 0xC4)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
LGRAY  = RGBColor(0xE2, 0xE8, 0xF0)
GRAY   = RGBColor(0x9C, 0xA3, 0xAF)
BGFILL = RGBColor(0xEE, 0xF2, 0xFF)

# ── Slide dimensions: 16:9, Keynote-compatible ─────────────────────────────────
W = Inches(13.333)
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

# Fix sldSz type: python-pptx defaults to "screen4x3" even for 16:9 dimensions.
# Keynote rejects the dimension/type mismatch — set the correct type.
prs_xml = prs.element
sldSz = prs_xml.find(qn('p:sldSz'))
if sldSz is not None:
    sldSz.set('type', 'screen16x9')

BLANK = prs.slide_layouts[6]   # truly blank layout


# ── Helpers ────────────────────────────────────────────────────────────────────

def add_rect(slide, left, top, width, height, fill_rgb=None, line_rgb=None, line_pt=0):
    shape = slide.shapes.add_shape(1, left, top, width, height)  # MSO_SHAPE_TYPE.RECTANGLE
    if fill_rgb:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_rgb
    else:
        shape.fill.background()
    if line_rgb and line_pt:
        shape.line.color.rgb = line_rgb
        shape.line.width = Pt(line_pt)
    else:
        shape.line.fill.background()
    return shape


def add_textbox(slide, text, left, top, width, height,
                font_name="Helvetica Neue", font_pt=22, bold=False, italic=False,
                color=NAVY, align=PP_ALIGN.LEFT, word_wrap=True):
    txb = slide.shapes.add_textbox(left, top, width, height)
    txb.word_wrap = word_wrap
    tf  = txb.text_frame
    tf.word_wrap = word_wrap
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    f = run.font
    f.name  = font_name
    f.size  = Pt(font_pt)
    f.bold  = bold
    f.italic = italic
    f.color.rgb = color
    return txb


def title_bar(slide, title_text, bar_h=Inches(0.9)):
    """Navy title bar across full width."""
    bar = add_rect(slide, 0, 0, W, bar_h, fill_rgb=NAVY)
    # gold accent line below bar
    add_rect(slide, 0, bar_h, W, Pt(4), fill_rgb=GOLD)
    add_textbox(slide, title_text,
                left=Inches(0.5), top=Inches(0.12),
                width=W - Inches(1), height=bar_h - Inches(0.12),
                font_pt=28, bold=True, color=WHITE)


def bullet_slide(slide, title, bullets, notes_text=""):
    """Standard content slide: title bar + bullet list."""
    title_bar(slide, title)
    y = Inches(1.1)
    body_h = H - y - Inches(0.3)
    txb = slide.shapes.add_textbox(Inches(0.55), y, W - Inches(0.9), body_h)
    txb.word_wrap = True
    tf = txb.text_frame
    tf.word_wrap = True
    for i, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_before = Pt(4)
        p.level = bullet.get("level", 0)
        run = p.add_run()
        run.text = bullet["text"]
        run.font.name = "Helvetica Neue"
        run.font.size  = Pt(bullet.get("size", 21))
        run.font.bold  = bullet.get("bold", False)
        run.font.color.rgb = bullet.get("color", NAVY)
    if notes_text:
        slide.notes_slide.notes_text_frame.text = notes_text


def figure_slide(slide, title, img_path, caption="", notes_text=""):
    """Full-width figure slide."""
    title_bar(slide, title)
    img = str(img_path) if img_path.exists() else None
    if img:
        pic = slide.shapes.add_picture(
            img, Inches(0.25), Inches(1.05),
            width=W - Inches(0.5)
        )
        # scale to fit height
        max_h = H - Inches(1.05) - (Inches(0.55) if caption else Inches(0.2))
        if pic.height > max_h:
            ratio = max_h / pic.height
            pic.height = int(pic.height * ratio)
            pic.width  = int(pic.width  * ratio)
        # center horizontally
        pic.left = int((W - pic.width) / 2)
    if caption:
        cap_top = H - Inches(0.55)
        add_textbox(slide, caption,
                    left=Inches(0.4), top=cap_top,
                    width=W - Inches(0.8), height=Inches(0.5),
                    font_pt=14, color=GRAY, italic=True)
    if notes_text:
        slide.notes_slide.notes_text_frame.text = notes_text


def title_slide(slide, title, subtitle, authors, affiliation, date):
    """Dark navy title slide."""
    # full navy background
    add_rect(slide, 0, 0, W, H, fill_rgb=NAVY)
    # gold accent bar at 60% height
    accent_y = Inches(4.8)
    add_rect(slide, Inches(0.8), accent_y, W - Inches(1.6), Pt(3), fill_rgb=GOLD)
    # title
    add_textbox(slide, title,
                left=Inches(0.8), top=Inches(1.4),
                width=W - Inches(1.6), height=Inches(1.5),
                font_pt=42, bold=True, color=WHITE)
    # subtitle
    add_textbox(slide, subtitle,
                left=Inches(0.8), top=Inches(2.85),
                width=W - Inches(1.6), height=Inches(1.1),
                font_pt=20, color=RGBColor(0xCB, 0xD5, 0xE1))
    # authors
    add_textbox(slide, authors,
                left=Inches(0.8), top=Inches(5.05),
                width=W - Inches(1.6), height=Inches(0.5),
                font_pt=19, bold=True, color=WHITE)
    # affiliation + date
    add_textbox(slide, f"{affiliation}  ·  {date}",
                left=Inches(0.8), top=Inches(5.55),
                width=W - Inches(1.6), height=Inches(0.4),
                font_pt=16, color=RGBColor(0x94, 0xA3, 0xB8))


def summary_slide(slide, points, closing, notes_text=""):
    """Dark navy summary slide."""
    add_rect(slide, 0, 0, W, H, fill_rgb=NAVY)
    add_textbox(slide, "Summary",
                left=Inches(0.8), top=Inches(0.3),
                width=W - Inches(1.6), height=Inches(0.7),
                font_pt=36, bold=True, color=WHITE)
    add_rect(slide, Inches(0.8), Inches(0.95), W - Inches(1.6), Pt(3), fill_rgb=GOLD)
    y = Inches(1.1)
    for pt in points:
        label_text = pt["label"]
        body_text  = pt["body"]
        add_textbox(slide, label_text,
                    left=Inches(0.8), top=y,
                    width=W - Inches(1.6), height=Inches(0.4),
                    font_pt=17, bold=True, color=GOLD)
        add_textbox(slide, body_text,
                    left=Inches(0.8), top=y + Inches(0.38),
                    width=W - Inches(1.6), height=Inches(0.5),
                    font_pt=20, color=WHITE)
        y += Inches(1.05)
    add_rect(slide, Inches(0.8), y, W - Inches(1.6), Pt(2), fill_rgb=GOLD)
    add_textbox(slide, closing,
                left=Inches(0.8), top=y + Inches(0.1),
                width=W - Inches(1.6), height=Inches(0.8),
                font_pt=18, italic=True, color=RGBColor(0xCB, 0xD5, 0xE1))
    if notes_text:
        slide.notes_slide.notes_text_frame.text = notes_text


def add_slide():
    return prs.slides.add_slide(BLANK)


# ── Slide 1: Title ─────────────────────────────────────────────────────────────
s = add_slide()
title_slide(s,
    title="Contact without Commitment",
    subtitle="Atomistic Characterisation of β-Lactoglobulin\nAdsorption Dynamics at the Air–Water Interface",
    authors="Chalakon Pornjariyawatch · Prapasiri Pongprayoon",
    affiliation="COMFHA, Kasetsart University",
    date="May 2026")
s.notes_slide.notes_text_frame.text = (
    "Welcome everyone. Today I'll be presenting our first paper on BLG adsorption at the air-water interface. "
    "This has been a long journey — I'll take you through the key findings in about 20 minutes."
)

# ── Slide 2: Today's Story ─────────────────────────────────────────────────────
s = add_slide()
bullet_slide(s, "Today's Story", [
    {"text": "THE PROBLEM", "bold": True, "color": GOLD},
    {"text": "BLG is the key foam stabiliser in milk — but how does it recognise the interface?", "level": 0},
    {"text": "Kinetic energy barrier measured for decades; mechanism unknown at atomic scale", "level": 0},
    {"text": "OUR APPROACH", "bold": True, "color": GOLD},
    {"text": "First unbiased atomistic MD of native BLG at the air–water interface — 4.00 µs total", "level": 0},
    {"text": "KEY FINDINGS", "bold": True, "color": GOLD},
    {"text": "Contact frequent (613 events); commitment absent (0 of 613)", "level": 0},
    {"text": "Calyx mobile; Loop CD/EF interface-induced; SASA ⊥ orientation", "level": 0},
],
notes_text="A quick roadmap before we dive in. I'll spend about 3 minutes on setup and approach, then the bulk on findings, then implications.")

# ── Slide 3: Why Milk Foam? ────────────────────────────────────────────────────
s = add_slide()
bullet_slide(s, "Why Milk Foam?", [
    {"text": "β-Lactoglobulin (BLG): dominant whey protein in bovine milk (~ 3 g/L)", "bold": False},
    {"text": "Principal stabiliser of milk foam — forms viscoelastic film at the AWI", "bold": False},
    {"text": " "},
    {"text": "Adsorption is slow: seconds to minutes from bulk solution", "bold": False},
    {"text": "A kinetic energy barrier distinguishes BLG from more flexible proteins", "bold": False},
    {"text": "Classical model (Graham & Phillips 1979): surface tension → global unfolding", "bold": False},
    {"text": "→ Assumed but never directly observed at atomic resolution", "color": BLUE},
],
notes_text=(
    "The motivation is the foam stability problem. BLG is the main protein in milk that makes foam stable. "
    "It adsorbs slowly — there's a measured kinetic barrier. The classical picture says global unfolding drives adsorption. "
    "But nobody has ever watched this at the atom level."
))

# ── Slide 4: The Atomistic Gap ─────────────────────────────────────────────────
s = add_slide()
bullet_slide(s, "The Atomistic Gap", [
    {"text": "Experiments (tensiometry, ellipsometry, neutron reflectometry):", "bold": True, "color": NAVY},
    {"text": "→ Time-averaged surface properties. No individual molecular events.", "level": 0, "color": GRAY},
    {"text": " "},
    {"text": "Prior MD studies of BLG at interfaces:", "bold": True, "color": NAVY},
    {"text": "→ Oil-water only (not AWI)", "level": 0},
    {"text": "→ Pre-positioned near interface — not starting from bulk", "level": 0},
    {"text": "→ ≤ 100 ns — far below the seconds-to-minutes adsorption timescale", "level": 0},
    {"text": " "},
    {"text": "This work: first unbiased atomistic simulation of native BLG at AWI", "bold": True, "color": GOLD},
],
notes_text=(
    "There's a gap between what experiments can tell us and what MD has addressed so far. "
    "Experiments give us time-averaged surface properties. Prior MD used the wrong interface or pre-positioned the protein. "
    "We fill this gap."
))

# ── Slide 5: System ────────────────────────────────────────────────────────────
s = add_slide()
figure_slide(s, "System: 12×12×35 nm Slab — 4.00 µs Cumulative",
    FIGS / "PAPER_FIG2_CONTACT_AB.png",
    caption="SET 1A: CENTER bulk-start (1000 ns) · SET 1B: R1, R2, R3 near-interface replicas (1000 ns each) · CHARMM36m + TIP3P",
    notes_text=(
        "Our system: a 12 by 12 by 35 nm slab with about 7 nm of water and vacuum on each side creating two air-water interfaces. "
        "One CENTER simulation starting from bulk, plus three replicas starting near the interface. "
        "Total: 4 microseconds of trajectory."
    ))

# ── Slide 6: Contact Metric ────────────────────────────────────────────────────
s = add_slide()
bullet_slide(s, "Contact Metric: Why Nearest-Atom Distance", [
    {"text": "BLG monomer is ~ 4 nm in diameter", "bold": False},
    {"text": "Centre-of-mass: stays 2–3 nm from interface even while touching it", "color": GRAY},
    {"text": "→ CoM distance is BLIND to actual surface contact", "bold": True, "color": RGBColor(0xDC,0x26,0x26)},
    {"text": " "},
    {"text": "Nearest-atom ≤ 0.3 nm: detects real surface touch", "bold": True, "color": GOLD},
    {"text": "→ Recovers 613 contact events vs. ~ 0 with CoM threshold", "color": NAVY},
    {"text": " "},
    {"text": "Penetration depth resolved to 0.01 nm precision", "bold": False},
],
notes_text=(
    "This is a key methodological point. The protein is 4 nm wide. "
    "Its centre of mass stays far from the interface even when an atom is touching it. "
    "Switching to the nearest atom threshold reveals that BLG visits the interface constantly."
))

# ── Slide 7: Contact Frequent — AB ────────────────────────────────────────────
s = add_slide()
figure_slide(s, "Contact is Frequent — CENTER + R1",
    FIGS / "PAPER_FIG2_CONTACT_AB.png",
    caption="CENTER: 97 events (12.5% of frames) · R1: 215 events (23.4% of frames) · Penetration up to 0.71 nm past the interface",
    notes_text=(
        "Here's what 4 microseconds looks like. The protein is constantly visiting the interface — 97 events in CENTER, 215 in R1. "
        "This is 12 to 23 percent of all frames. And single atoms penetrate up to 0.71 nm past the interface plane."
    ))

# ── Slide 8: Commitment Rare — CD ─────────────────────────────────────────────
s = add_slide()
figure_slide(s, "Commitment is Rare — R2 + R3",
    FIGS / "PAPER_FIG2_CONTACT_CD.png",
    caption="R2: 156 events · R3: 145 events · Total 613 — only 6 sustain ≥ 10 ns · None commits",
    notes_text=(
        "Same story in R2 and R3. 613 total contact events across 4 µs. "
        "But 607 of them terminate within 10 nanoseconds. "
        "And none — not a single one — commits to stable adsorption. "
        "That's the contact-commitment dichotomy."
    ))

# ── Slide 9: Long Events Table ─────────────────────────────────────────────────
s = add_slide()
title_bar(s, "The 6 Long-Residency Events: All Non-Activated")
# Draw simple table manually
table_data = [
    ["Replica", "Duration", "Event-mean SASA", "Angle", "Committed?"],
    ["CENTER",  "~10 ns",  "28.8 nm²",  "~52°", "No"],
    ["R1",      "59 ns",   "29.1 nm²",  "~63°", "No"],
    ["R1",      "~12 ns",  "28.5 nm²",  "~68°", "No"],
    ["R1",      "~11 ns",  "30.5 nm²",  "~44°", "No"],
    ["R1",      "~10 ns",  "28.9 nm²",  "~75°", "No"],
    ["R3",      "~10 ns",  "30.1 nm²",  "~56°", "No"],
]
col_widths = [Inches(1.6), Inches(1.5), Inches(2.3), Inches(1.5), Inches(2.0)]
row_h = Inches(0.48)
tbl_top = Inches(1.05)
x_start = Inches(0.55)
for ri, row in enumerate(table_data):
    x = x_start
    for ci, cell_text in enumerate(row):
        is_header = (ri == 0)
        bg = NAVY if is_header else (BGFILL if ri % 2 == 0 else WHITE)
        fg = WHITE if is_header else (GOLD if ci == 1 and ri > 0 and "59" in cell_text else NAVY)
        rect = add_rect(s, x, tbl_top + ri * row_h,
                        col_widths[ci], row_h,
                        fill_rgb=bg, line_rgb=LGRAY, line_pt=0.5)
        add_textbox(s, cell_text,
                    left=x + Pt(6), top=tbl_top + ri * row_h + Pt(6),
                    width=col_widths[ci] - Pt(12), height=row_h - Pt(6),
                    font_pt=17, bold=is_header, color=fg)
        x += col_widths[ci]
add_textbox(s,
    "All 6 events: SASA 28.5–30.5 nm² · zero activation-criterion frames in every event\n"
    "The protein sustains 59 ns of contact — still doesn't commit",
    left=Inches(0.55), top=Inches(4.55),
    width=W - Inches(0.9), height=Inches(0.9),
    font_pt=19, italic=True, color=BLUE)
s.notes_slide.notes_text_frame.text = (
    "Let me show you the six long events explicitly. The longest — 59 nanoseconds in R1 — is remarkable. "
    "That's not a brief fluctuation. But the SASA stays around 29 nm² throughout and never reaches 32.1. "
    "Contact is not commitment."
)

# ── Slide 10: Global Compactness ───────────────────────────────────────────────
s = add_slide()
figure_slide(s, "Global Fold Preserved — Never Unfolds in 4 µs",
    FIGS / "PAPER_FIG3_ACTIVATION.png",
    caption="Rg = 1.496 ± 0.009 nm (R1) — no trend · β-barrel RMSD 0.21–0.24 nm · α-helix RMSD ≤ 0.137 nm",
    notes_text=(
        "Now let's look at what's happening structurally. The radius of gyration is rock solid — 1.496 nm, no trend. "
        "The beta-barrel stays assembled. The alpha-helix intact. "
        "Under 4 microseconds of unbiased dynamics, the protein never globally unfolds. "
        "The Graham and Phillips surface-denaturation picture does not describe what we observe."
    ))

# ── Slide 11: Calyx Breathes ───────────────────────────────────────────────────
s = add_slide()
figure_slide(s, "Calyx Breathes — Stationary Stochastic Process",
    FIGS / "PAPER_FIG3_ACTIVATION.png",
    caption="SASA: 24–37 nm² (PBC-corrected) · Recurring fluctuations ~30–40 ns · Not a one-shot pre-adsorption event",
    notes_text=(
        "But even though the protein doesn't globally unfold, the calyx is mobile. "
        "SASA fluctuates between 24 and 37 nm² with a period of roughly 30 to 40 ns. "
        "This isn't a one-shot opening before adsorption — it's a stationary stochastic process happening constantly."
    ))

# ── Slide 12: RMSF Loop Shift ──────────────────────────────────────────────────
s = add_slide()
figure_slide(s, "Interface Shifts Which Loop Dominates",
    FIGS / "PAPER_FIG3_ACTIVATION.png",
    caption="Bulk (CENTER): Loop BC dominant, RMSF 0.54 nm · Near-AWI (R1): Loop CD/EF rises to 0.39 nm — interface-induced conformational preference",
    notes_text=(
        "This is one of the more striking findings. In bulk, Loop BC is the most flexible region. "
        "Near the air-water interface, Loop CD/EF — which sits directly above the hydrophobic calyx — rises to become the dominant flexible region. "
        "This is an interface-induced conformational preference. Activation is loop-mediated, not global."
    ))

# ── Slide 13: SASA Distribution ───────────────────────────────────────────────
s = add_slide()
figure_slide(s, "SASA Distribution: All Four Replicas",
    FIGS / "Fig4_optionA.png",
    caption="All replicas: SASA 24–37 nm² · p95 = 32.1 nm² (distribution-based) · No replica reaches a distinct activated regime",
    notes_text=(
        "Figure 4 shows two panels. The left panel is the KDE of SASA for all four replicas. "
        "They're remarkably consistent — all confined to 24 to 37 nm². "
        "We define p95 at 32.1 nm² as a distribution-based threshold. No replica crosses it cleanly into an activated state."
    ))

# ── Slide 14: Orientation Independent ─────────────────────────────────────────
s = add_slide()
figure_slide(s, "Orientation: Uniform and Independent of SASA",
    FIGS / "Fig4_optionA.png",
    caption="Pearson r = +0.006 · No quadrant clustering · Obs/Indep ≈ 1.0 across SASA threshold sweep (27–33 nm²)",
    notes_text=(
        "The right panel is the 2D joint distribution of calyx angle versus SASA. "
        "If orientation and SASA were coupled — as a two-factor gate model would predict — you'd see density concentrated in one quadrant. "
        "Instead it's uniform. Pearson r is plus 0.006. "
        "We tested this across every reasonable threshold and the result holds: the two coordinates are independent."
    ))

# ── Slide 15: Block Bootstrap ──────────────────────────────────────────────────
s = add_slide()
bullet_slide(s, "Statistical Robustness: Block Bootstrap", [
    {"text": "Challenge: SASA autocorrelation = 232 ns", "bold": True, "color": NAVY},
    {"text": "→ i.i.d. assumption breaks down for standard confidence intervals", "color": GRAY},
    {"text": " "},
    {"text": "Block bootstrap approach:", "bold": True, "color": NAVY},
    {"text": "→ Autocorrelation measured: 81–394 ns per replica", "level": 0},
    {"text": "→ Block length: 232 ns (conservative)", "level": 0},
    {"text": "→ Effective N ≈ 17 independent observations across 4 µs", "level": 0, "bold": True},
    {"text": " "},
    {"text": "Block bootstrap 95% CI: [−0.09, +0.11]", "bold": True, "color": GOLD},
    {"text": "→ Rules out |r| > 0.11 at 95% confidence", "bold": True, "color": GOLD},
],
notes_text=(
    "I want to spend a moment on statistics because it matters here. "
    "SASA has an autocorrelation time of 232 nanoseconds. "
    "So the effective number of independent observations across 4 microseconds is only about 17. "
    "We account for this with block bootstrap. The 95% confidence interval is minus 0.09 to plus 0.11. "
    "We can confidently rule out any linear coupling stronger than about 0.11."
))

# ── Slide 16: PBC Lesson ───────────────────────────────────────────────────────
s = add_slide()
bullet_slide(s, "Scientific Integrity: The PBC Lesson", [
    {"text": "During analysis we found a periodic-boundary artefact:", "bold": True, "color": NAVY},
    {"text": "freeSASA without MDAnalysis unwrap → split atoms across PBC boundaries", "level": 0, "color": GRAY},
    {"text": " "},
    {"text": "Before fix:  SASA 45–62 nm²  ·  21 gate-open events", "bold": True, "color": RGBColor(0xDC,0x26,0x26)},
    {"text": "After  fix:  SASA 24–37 nm²  ·  0 gate-open events", "bold": True, "color": GOLD},
    {"text": " "},
    {"text": "Fix: mda_unwrap transformation in MDAnalysis — applied to all 8006 gate frames", "level": 0},
    {"text": "Full disclosure in Methods section of the paper", "level": 0},
    {"text": " "},
    {"text": "All results in this talk use PBC-corrected data", "bold": True, "color": NAVY},
],
notes_text=(
    "I need to mention this because scientific integrity matters. "
    "Midway through the project we found that freeSASA without unwrapping the trajectory inflated SASA to 45 to 62 nm². "
    "After the fix, it's 24 to 37. The gate-open events went from 21 to zero. "
    "We disclose this fully in the Methods. Everything you've seen today uses the corrected analysis."
))

# ── Slide 17: What This Means ──────────────────────────────────────────────────
s = add_slide()
bullet_slide(s, "What This Means", [
    {"text": "Prior models assumed:", "bold": True, "color": NAVY},
    {"text": "Contact → immediate unfolding → adsorption (Graham & Phillips)", "level": 0, "color": GRAY},
    {"text": "OR: brief contacts don't matter", "level": 0, "color": GRAY},
    {"text": " "},
    {"text": "What we observe:", "bold": True, "color": NAVY},
    {"text": "59 ns contact without commitment → genuine kinetic bottleneck, not sampling artefact", "level": 0},
    {"text": "r = +0.006 → no simple two-factor coupling as commitment trigger", "level": 0},
    {"text": "Global fold preserved → commitment mechanism is not global unfolding", "level": 0},
    {"text": " "},
    {"text": "First characterisation of the pre-commitment contact ensemble", "bold": True, "color": GOLD},
    {"text": "→ Required foundation for enhanced-sampling calculations", "level": 0, "color": GOLD},
],
notes_text=(
    "So what does this all add up to? The prior models don't fit. "
    "59 nanoseconds of contact without commitment is a real kinetic bottleneck. "
    "SASA and orientation are decoupled — there's no simple two-factor gate. "
    "And we never see global unfolding. "
    "What we have is the first characterisation of the pre-commitment state — and that's what you need before you can design the enhanced sampling."
))

# ── Slide 18: Implications ─────────────────────────────────────────────────────
s = add_slide()
bullet_slide(s, "Implications", [
    {"text": "Loop CD/EF: the engineering target", "bold": True, "color": GOLD},
    {"text": "Becomes dominant near AWI — interface-induced conformational preference", "level": 0},
    {"text": "Mutations increasing calyx flexibility likely more effective than bulk hydrophobicity", "level": 0},
    {"text": " "},
    {"text": "Dimer complexity (BLG dominant as dimer above ~50 µM)", "bold": True, "color": NAVY},
    {"text": "Steric: dimer interface occludes Loop CD/EF region", "level": 0},
    {"text": "Kinetic: dimer must dissociate before calyx presents unobstructed", "level": 0},
    {"text": "Orientational: oblate shape alters rotational diffusion and θ distribution", "level": 0},
    {"text": " "},
    {"text": "Lipocalin family: quantitative SASA/orientation baseline for mutant comparison", "level": 0},
],
notes_text=(
    "Two practical implications. First, Loop CD/EF is the engineering target for foam stabilisers. "
    "If you want to mutate BLG to adsorb faster, target calyx flexibility, not just surface hydrophobicity. "
    "Second, we simulated the monomer. The native dimer adds complexity — steric, kinetic, orientational — "
    "that our next simulations will address."
))

# ── Slide 19: What's Next ──────────────────────────────────────────────────────
s = add_slide()
bullet_slide(s, "What's Next", [
    {"text": "Paper 1 (this work) — submission pipeline", "bold": True, "color": GOLD},
    {"text": "Zenodo upload → DOI → JCIS submission (IF ~9)", "level": 0},
    {"text": " "},
    {"text": "Enhanced sampling — the clear next step", "bold": True, "color": NAVY},
    {"text": "Metadynamics / umbrella sampling along (SASA, θ)", "level": 0},
    {"text": "SET 1D data available as baseline starting conformations", "level": 0},
    {"text": " "},
    {"text": "Paper 2 — β-Casein at AWI", "bold": True, "color": NAVY},
    {"text": "AlphaFold2 structure ready; comparison with BLG pre-commitment ensemble", "level": 0},
    {"text": " "},
    {"text": "BLG dimer + TIP4P/2005 cross-check", "level": 0},
],
notes_text=(
    "The immediate next step is getting Paper 1 submitted to JCIS — we're waiting on Zenodo DOI and P.P. sign-off. "
    "The scientific next step is enhanced sampling to characterise the commitment mechanism. "
    "And Paper 2 on beta-casein is already started — it'll be an interesting comparison."
))

# ── Slide 20: Summary ──────────────────────────────────────────────────────────
s = add_slide()
summary_slide(s,
    points=[
        {"label": "1. Contact without commitment",
         "body":  "613 events across 4 µs — only 6 sustain ≥ 10 ns, none commits to adsorption"},
        {"label": "2. Compact globally, loop-mediated locally",
         "body":  "Rg stable; β-barrel intact; Loop CD/EF mediates calyx exposure near AWI"},
        {"label": "3. SASA ⊥ orientation on the µs timescale",
         "body":  "r = +0.006, CI [−0.09, +0.11] — rules out |r| > 0.11 coupling"},
    ],
    closing="4 µs of unbiased MD resolves the pre-commitment contact ensemble.\nCommitment requires enhanced sampling to characterise.\n\nThank you — questions?",
    notes_text=(
        "Three takeaways. Contact frequent, commitment absent. Global fold preserved, loop-mediated calyx mobility. "
        "SASA and orientation statistically independent. "
        "These together define what unbiased microsecond MD can and cannot tell us about BLG adsorption. "
        "Thank you."
    ))

# ── Save ───────────────────────────────────────────────────────────────────────
prs.save(str(OUT))
print(f"Saved: {OUT}  ({OUT.stat().st_size // 1024} KB)")
print(f"Slides: {len(prs.slides)}")
print()
print("To open in Keynote on Mac:")
print("  1. Copy presentation.pptx to macOS")
print("  2. Open in Keynote → it imports automatically")
print("  3. File → Save As → .key")
