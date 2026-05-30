"""
build_presentation.py  v2
=========================
Redesigned: action titles (Gabberflast/Naegle principle), section colour coding,
conclusions slide that stays on screen during Q&A, panel labels on multi-panel slides.

Keynote-compatible: Marp PPTX base → strip image slides → inject live text + figures
→ deduplicate ZIP → push to Mac.
"""
import zipfile, subprocess, shutil
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

ROOT = Path(__file__).resolve().parent.parent
FIGS = ROOT / "slides" / "figures"
BASE = ROOT / "slides" / "slides_base.pptx"
OUT  = ROOT / "slides" / "presentation.pptx"
MAC_HOST = "cp@100.118.196.110"
MAC_PATH = "/Users/cp/Documents/CP/COMFHA_Talk_May2026/presentation.pptx"

# ── Colour palette ─────────────────────────────────────────────────────────────
NAVY    = RGBColor(0x1A, 0x27, 0x44)   # primary brand
GOLD    = RGBColor(0xC9, 0xA8, 0x4C)   # highlights
BLUE    = RGBColor(0x44, 0x72, 0xC4)   # secondary
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
LGRAY   = RGBColor(0xF1, 0xF5, 0xF9)   # light panel background
GRAY    = RGBColor(0x9C, 0xA3, 0xAF)
RED     = RGBColor(0xDC, 0x26, 0x26)
GREEN   = RGBColor(0x16, 0xA3, 0x4A)

# Section accent colours
PURPLE  = RGBColor(0x7C, 0x3A, 0xED)   # problem / motivation
OCEAN   = RGBColor(0x03, 0x69, 0xA1)   # method / approach
EMERALD = RGBColor(0x04, 0x78, 0x57)   # results
CRIMSON = RGBColor(0xDC, 0x26, 0x26)   # integrity / correction
AMBER   = RGBColor(0xB4, 0x53, 0x09)   # implications

# Slide geometry (16:9, 13.333 × 7.5 in)
W   = Inches(13.333); H   = Inches(7.5)
BAR = Inches(1.0);    MAR = Inches(0.55); BW = W - MAR * 2

# Section tag tuples (label_text, colour) for bar()
T_PROB = ("PROBLEM",      PURPLE)
T_METH = ("METHOD",       OCEAN)
T_RESL = ("RESULTS",      EMERALD)
T_INTG = ("INTEGRITY",    CRIMSON)
T_IMPL = ("IMPLICATIONS", AMBER)
T_NEXT = ("NEXT",         BLUE)

# ── Presentation base ──────────────────────────────────────────────────────────
prs = Presentation(str(BASE))
prs.slide_width = W; prs.slide_height = H
sldIdLst = prs.slides._sldIdLst
for el in list(sldIdLst): sldIdLst.remove(el)
LAYOUT = prs.slide_layouts[0]

def S(): return prs.slides.add_slide(LAYOUT)

# ── Low-level drawing helpers ──────────────────────────────────────────────────

def rect(slide, l, t, w, h, fill):
    sp = slide.shapes.add_shape(1, l, t, w, h)
    sp.fill.solid(); sp.fill.fore_color.rgb = fill; sp.line.fill.background()
    return sp

def txt(slide, text, l, t, w, h, size=20, color=NAVY, bold=False,
        align=PP_ALIGN.LEFT, wrap=True):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tb.word_wrap = wrap; tf = tb.text_frame; tf.word_wrap = wrap
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.color.rgb = color; r.font.bold = bold
    r.font.name = "Helvetica Neue Bold" if bold else "Helvetica Neue"
    return tb

# ── Structural components ──────────────────────────────────────────────────────

def bar(slide, title, accent=GOLD, section_tag=None):
    """Navy title bar with coloured bottom accent stripe + optional section pill."""
    rect(slide, 0, 0, W, BAR, NAVY)
    rect(slide, 0, BAR, W, Pt(4), accent)          # section-coloured stripe

    tag_reserve = Inches(1.35) if section_tag else Inches(0.1)
    tb = slide.shapes.add_textbox(MAR, Inches(0.11), BW - tag_reserve, BAR - Inches(0.19))
    tf = tb.text_frame; tf.word_wrap = True
    r = tf.paragraphs[0].add_run()
    r.text = title; r.font.size = Pt(22); r.font.bold = True
    r.font.color.rgb = WHITE; r.font.name = "Helvetica Neue Bold"

    if section_tag:
        tag_text, tag_color = section_tag
        pw = Inches(1.22); ph = Inches(0.29)
        pl = W - MAR - pw; pt_y = Inches(0.145)
        rect(slide, pl, pt_y, pw, ph, tag_color)
        tb2 = slide.shapes.add_textbox(pl, pt_y, pw, ph)
        p2 = tb2.text_frame.paragraphs[0]; p2.alignment = PP_ALIGN.CENTER
        r2 = p2.add_run(); r2.text = tag_text
        r2.font.size = Pt(9); r2.font.bold = True
        r2.font.color.rgb = WHITE; r2.font.name = "Helvetica Neue Bold"

def dark_bg(slide, color=NAVY):
    rect(slide, 0, 0, W, H, color)

def slide_num(slide, n, total=20):
    tb = slide.shapes.add_textbox(W - Inches(1.1), H - Inches(0.32), Inches(1.0), Inches(0.28))
    r = tb.text_frame.paragraphs[0].add_run()
    r.text = f"{n}/{total}"; r.font.size = Pt(12); r.font.color.rgb = GRAY
    r.font.name = "Helvetica Neue"
    tb.text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT

def fig(slide, path, caption=None, panel_label=None):
    """Embed figure with optional caption and coloured panel badge."""
    top = BAR + Pt(12)
    max_h = H - top - (Inches(0.46) if caption else Inches(0.18))
    if path.exists():
        pic = slide.shapes.add_picture(str(path), Inches(0.15), top, width=W - Inches(0.3))
        if pic.height > max_h:
            ratio = max_h / pic.height
            pic.height = int(pic.height * ratio); pic.width = int(pic.width * ratio)
        pic.left = int((W - pic.width) / 2)
    else:
        txt(slide, f"[ {path.name} ]", MAR, top, BW, max_h, size=18, color=GRAY,
            align=PP_ALIGN.CENTER)
    if caption:
        txt(slide, caption, MAR, H - Inches(0.44), BW, Inches(0.40), size=12, color=GRAY)
    if panel_label:
        bw_p = Inches(0.50); bh_p = Inches(0.26)
        rect(slide, Inches(0.20), top + Inches(0.05), bw_p, bh_p, GOLD)
        tb_p = slide.shapes.add_textbox(Inches(0.20), top + Inches(0.05), bw_p, bh_p)
        p_p = tb_p.text_frame.paragraphs[0]; p_p.alignment = PP_ALIGN.CENTER
        r_p = p_p.add_run(); r_p.text = panel_label
        r_p.font.size = Pt(11); r_p.font.bold = True
        r_p.font.color.rgb = NAVY; r_p.font.name = "Helvetica Neue Bold"

def buls(slide, items, extra=Inches(0)):
    """
    Bullet list. items = list of (text, bold, color) or (text, bold, color, indent).
    indent=1 gives secondary indentation with – bullet.
    """
    top = BAR + Pt(16)
    h = H - top - Inches(0.38) - extra
    tb = slide.shapes.add_textbox(MAR, top, BW, h)
    tb.word_wrap = True; tf = tb.text_frame; tf.word_wrap = True
    first = True
    for item in items:
        text, bold, color = item[:3]
        indent = item[3] if len(item) == 4 else 0
        p = tf.paragraphs[0] if first else tf.add_paragraph(); first = False
        p.space_before = Pt(2 if indent else 4)
        r = p.add_run()
        if bold:
            bullet = "▸ "        # ▸
        elif indent:
            bullet = "    – "    # em-indent + –
        else:
            bullet = "• "        # •
        r.text = bullet + text
        r.font.size = Pt(18 if bold else 17)
        r.font.color.rgb = color; r.font.bold = bold
        r.font.name = "Helvetica Neue Bold" if bold else "Helvetica Neue"

def callout(slide, text, accent_color=GOLD):
    """Bottom callout box with coloured left border."""
    y = H - Inches(0.96); box_h = Inches(0.62)
    rect(slide, MAR - Inches(0.05), y, BW + Inches(0.10), box_h, LGRAY)
    rect(slide, MAR - Inches(0.05), y, Pt(5), box_h, accent_color)
    txt(slide, text, MAR + Pt(10), y + Inches(0.10), BW - Pt(10), box_h - Inches(0.10),
        size=16, color=NAVY, align=PP_ALIGN.CENTER)

# =============================================================================
# SLIDES 1 – 20
# =============================================================================

# ── S1  Title ─────────────────────────────────────────────────────────────────
s = S(); dark_bg(s)
# Subtle dark corner accent (atmosphere, not decoration)
rect(s, W - Inches(5.0), H - Inches(2.5), Inches(5.0), Inches(2.5), RGBColor(0x0D, 0x1B, 0x36))
# Gold separator line under main title
rect(s, MAR, Inches(2.72), Inches(6.0), Pt(3), GOLD)
txt(s, "Contact without Commitment", MAR, Inches(0.70), BW, Inches(1.90),
    size=50, color=WHITE, bold=True)
txt(s, "Atomistic Characterisation of β-Lactoglobulin Adsorption Dynamics"
       "\nat the Air–Water Interface",
    MAR, Inches(2.80), BW, Inches(1.10), size=19, color=RGBColor(0xCB, 0xD5, 0xE1))
txt(s, "Chalakon Pornjariyawatch  |  Prapasiri Pongprayoon",
    MAR, Inches(4.05), BW, Inches(0.50), size=19, color=WHITE, bold=True)
txt(s, "COMFHA — Kasetsart University  |  Lab Group Seminar  |  May 2026",
    MAR, Inches(4.54), BW, Inches(0.40), size=14, color=GRAY)
# Target journal badge
rect(s, MAR, Inches(5.30), Inches(5.2), Inches(0.46), RGBColor(0x25, 0x38, 0x60))
txt(s, "Submitting → Journal of Colloid and Interface Science  (IF ∼9)",
    MAR + Pt(12), Inches(5.34), Inches(5.0), Inches(0.38), size=13, color=GOLD)
slide_num(s, 1)

# ── S2  Overview ──────────────────────────────────────────────────────────────
s = S()
bar(s, "BLG Visits the AWI Constantly — but Never Commits to Adsorption", GOLD)
buls(s, [
    ("THE PROBLEM",                                                        True,  PURPLE),
    ("BLG stabilises milk foam — adsorbs over seconds, mechanism unknown atomically", False, NAVY, 1),
    ("Prior MD: oil-water only, ≤10 ns, pre-positioned — never native BLG at AWI", False, NAVY, 1),
    ("OUR APPROACH",                                                        True,  OCEAN),
    ("First unbiased atomistic MD of native BLG at AWI — 4.00 µs total",          False, NAVY, 1),
    ("KEY FINDINGS",                                                        True,  EMERALD),
    ("613 contact events across 4 replicas — frequent, brief, non-committing",         False, NAVY, 1),
    ("Calyx SASA ⊥ orientation (r = +0.006) — no simple two-factor gate",          False, NAVY, 1),
], extra=Inches(0.65))
callout(s, "Four acts:  Problem  →  Approach  →  Findings  →  Implications", GOLD)
slide_num(s, 2)

# ── S3  Why milk foam? ────────────────────────────────────────────────────────
s = S()
bar(s, "BLG Controls Milk Foam Stability — Its AWI Entry Mechanism Is Unknown",
    PURPLE, T_PROB)
buls(s, [
    ("β-Lactoglobulin (BLG): dominant whey protein in bovine milk (∼3 g/L)",    False, NAVY),
    ("Principal stabiliser of milk foam — forms viscoelastic film at the AWI",       False, NAVY),
    ("Adsorption is slow: seconds to minutes from bulk solution",                         False, NAVY),
    ("Kinetic barrier distinguishes BLG from more flexible, disordered proteins",         False, NAVY),
    ("Classical model (Graham & Phillips 1979): surface tension → global unfolding", False, GRAY),
    ("Assumed — never directly observed at atomic resolution",                       True,  CRIMSON),
], extra=Inches(0.70))
callout(s, "How does a folded, water-soluble protein recognise and populate the AWI?", PURPLE)
slide_num(s, 3)

# ── S4  Atomistic gap ─────────────────────────────────────────────────────────
s = S()
bar(s, "No Prior Atomistic MD Has Placed Native BLG at an AWI Unbiased", PURPLE, T_PROB)
buls(s, [
    ("Experiments (tensiometry, ellipsometry, neutron reflectometry):",                    True,  NAVY),
    ("Time-averaged surface properties. No individual molecular events.",                  False, GRAY, 1),
    ("Prior MD studies of BLG at interfaces:",                                             True,  NAVY),
    ("Oil-water interfaces only — not air-water",                                     False, NAVY, 1),
    ("Pre-positioned near interface — not starting from bulk",                        False, NAVY, 1),
    ("≤10–100 ns — far below seconds-to-minutes adsorption timescale",      False, NAVY, 1),
    ("This work: first unbiased atomistic MD of native BLG at an air-water interface",     True,  PURPLE),
])
slide_num(s, 4)

# ── S5  System ────────────────────────────────────────────────────────────────
s = S()
bar(s, "4.00 µs Unbiased MD in a 12×12×35 nm Slab with Two AWI Regions",
    OCEAN, T_METH)
fig(s, FIGS / "PAPER_FIG2_CONTACT_AB.png",
    "SET 1A: CENTER bulk-start (1000 ns)  |  SET 1B: R1, R2, R3 near-interface starts (1000 ns each)"
    "  |  CHARMM36m + TIP3P")
slide_num(s, 5)

# ── S6  Contact metric ────────────────────────────────────────────────────────
s = S()
bar(s, "Centre-of-Mass Distance Misses All Contact Events — Nearest-Atom ≤0.3 nm Does Not",
    OCEAN, T_METH)
buls(s, [
    ("BLG monomer is ∼4 nm in diameter",                                              False, NAVY),
    ("CoM stays 2–3 nm from interface even while touching it",                        False, GRAY),
    ("→ CoM threshold is BLIND to actual surface contact",                            True,  CRIMSON),
    ("Nearest-atom ≤0.3 nm: detects real surface touch",                              True,  GOLD),
    ("Recovers 613 contact events vs. ∼0 with CoM threshold",                         False, NAVY),
    ("Penetration depth resolved to 0.01 nm precision",                                    False, NAVY),
], extra=Inches(0.70))
callout(s, "Switching to nearest-atom distance reveals BLG visits the interface constantly", OCEAN)
slide_num(s, 6)

# ── S7  Contact frequent ──────────────────────────────────────────────────────
s = S()
bar(s, "BLG Contacts the AWI 312 Times Across CENTER + R1, Penetrating Up to 0.71 nm",
    EMERALD, T_RESL)
fig(s, FIGS / "PAPER_FIG2_CONTACT_AB.png",
    "CENTER: 97 events (12.5% of frames)  |  R1: 215 events (23.4%)  |  Max penetration 0.71 nm below interface plane")
slide_num(s, 7)

# ── S8  Commitment rare ───────────────────────────────────────────────────────
s = S()
bar(s, "All 613 Contact Events End Without Adsorption — Only 6 Last ≥10 ns",
    EMERALD, T_RESL)
fig(s, FIGS / "PAPER_FIG2_CONTACT_CD.png",
    "R2: 156 events  |  R3: 145 events  |  613 total across 4 replicas  |  None commits to adsorption")
slide_num(s, 8)

# ── S9  Long-event table ──────────────────────────────────────────────────────
s = S()
bar(s, "Even a 59 ns Contact Shows No Calyx Activation", EMERALD, T_RESL)
buls(s, [
    ("Replica  |  Duration  |  Mean SASA    |  Angle  |  Committed?",  True,  NAVY),
    ("CENTER   |  ∼10 ns   |  28.8 nm²  |  ∼52°  |  No",  False, NAVY),
    ("R1       |   59 ns    |  29.1 nm²  |  ∼63°  |  No   ← longest event", True, GOLD),
    ("R1       |  ∼12 ns   |  28.5 nm²  |  ∼68°  |  No",  False, NAVY),
    ("R1       |  ∼11 ns   |  30.5 nm²  |  ∼44°  |  No",  False, NAVY),
    ("R1       |  ∼10 ns   |  28.9 nm²  |  ∼75°  |  No",  False, NAVY),
    ("R3       |  ∼10 ns   |  30.1 nm²  |  ∼56°  |  No",  False, NAVY),
], extra=Inches(0.52))
txt(s, "All 6 events: SASA 28.5–30.5 nm²  |  angle 44–75°  |  zero activation frames  |  contact ≠ adsorption",
    MAR, H - Inches(0.86), BW, Inches(0.44), size=14, color=BLUE)
slide_num(s, 9)

# ── S10  Global compactness ───────────────────────────────────────────────────
s = S()
bar(s, "The Global β-Barrel Fold Remains Compact Across All 4 µs", EMERALD, T_RESL)
fig(s, FIGS / "PAPER_FIG3_ACTIVATION.png",
    "Rg = 1.496 ± 0.009 nm — flat  |  β-barrel RMSD 0.21–0.24 nm  |  α-helix intact  |  no trend in 4 µs",
    panel_label="(b) Rg")
slide_num(s, 10)

# ── S11  Calyx breathes ───────────────────────────────────────────────────────
s = S()
bar(s, "Calyx SASA Fluctuates 24–37 nm² — Not a One-Shot Pre-Adsorption Event",
    EMERALD, T_RESL)
fig(s, FIGS / "PAPER_FIG3_ACTIVATION.png",
    "SASA 24–37 nm² (PBC-corrected)  |  Recurring fluctuations ∼30–40 ns  |  Stationary stochastic process",
    panel_label="(a) SASA")
slide_num(s, 11)

# ── S12  Loop shift ───────────────────────────────────────────────────────────
s = S()
bar(s, "The AWI Selectively Amplifies Loop CD/EF Mobility Over Loop BC",
    EMERALD, T_RESL)
fig(s, FIGS / "PAPER_FIG3_ACTIVATION.png",
    "Bulk: Loop BC dominant (RMSF 0.54 nm)  |  Near AWI: Loop CD/EF rises to 0.39 nm  |  Interface-induced shift",
    panel_label="(c) RMSF")
slide_num(s, 12)

# ── S13  SASA distribution ────────────────────────────────────────────────────
s = S()
bar(s, "All Four Replicas Share the Same SASA Distribution — No Activated Subpopulation",
    EMERALD, T_RESL)
fig(s, FIGS / "Fig4_optionA.png",
    "All replicas: SASA 24–37 nm²  |  p95 = 32.1 nm²  |  No distinct high-SASA activated regime",
    panel_label="(a) KDE")
slide_num(s, 13)

# ── S14  Orientation ──────────────────────────────────────────────────────────
s = S()
bar(s, "Calyx Exposure and Orientation Are Statistically Independent (r = +0.006)",
    EMERALD, T_RESL)
fig(s, FIGS / "Fig4_optionA.png",
    "Pearson r = +0.006  |  No quadrant clustering  |  Obs/Indep ≈1.0 across threshold sweep",
    panel_label="(b) 2D")
slide_num(s, 14)

# ── S15  Block bootstrap ──────────────────────────────────────────────────────
s = S()
bar(s, "Block Bootstrap Rules Out |r| > 0.11 at 95% Confidence (N_eff ≈17)",
    EMERALD, T_RESL)
buls(s, [
    ("Challenge: SASA autocorrelation = 232 ns per replica",                               True,  NAVY),
    ("i.i.d. assumption breaks — standard CIs vastly underestimate uncertainty",      False, GRAY, 1),
    ("Block bootstrap solution:",                                                           True,  OCEAN),
    ("Autocorrelation range: 81–394 ns per replica",                                  False, NAVY, 1),
    ("Block length: 232 ns (average across replicas)",                                     False, NAVY, 1),
    ("Effective N ≈17 independent observations across 4 µs",                      False, NAVY, 1),
    ("Block bootstrap 95% CI: [−0.09, +0.11]",                                        True,  EMERALD),
], extra=Inches(0.70))
callout(s, "SASA and calyx orientation are statistically independent on the µs timescale", EMERALD)
slide_num(s, 15)

# ── S16  PBC integrity ────────────────────────────────────────────────────────
s = S()
bar(s, "A PBC Artefact Inflated SASA by 2× — Found, Corrected, and Fully Disclosed",
    CRIMSON, T_INTG)
buls(s, [
    ("What happened:",                                                                      True,  NAVY),
    ("freeSASA without MDAnalysis unwrap → atoms split across PBC boundaries",        False, GRAY, 1),
    ("Before fix: SASA 45–62 nm²  |  21 gate-open events",                       True,  CRIMSON),
    ("After  fix: SASA 24–37 nm²  |   0 gate-open events",                       True,  GREEN),
    ("Fix applied: mda_unwrap on all 8006 gate frames",                                    False, NAVY),
    ("Full disclosure in the Methods section of the paper",                                False, NAVY),
    ("All results in this talk use PBC-corrected data",                                    True,  NAVY),
])
slide_num(s, 16)

# ── S17  What this means ──────────────────────────────────────────────────────
s = S()
bar(s, "Commitment Requires a Mechanism Beyond Calyx–Orientation Coupling",
    AMBER, T_IMPL)
buls(s, [
    ("Prior models assumed:",                                                               True,  NAVY),
    ("Contact → global unfolding → adsorption (Graham & Phillips 1979)",         False, GRAY, 1),
    ("Or: brief contacts do not matter to adsorption",                                     False, GRAY, 1),
    ("What we observe instead:",                                                            True,  AMBER),
    ("59 ns contact without commitment → genuine kinetic bottleneck",                 False, NAVY, 1),
    ("r = +0.006 → no simple calyx–orientation coupling as trigger",             False, NAVY, 1),
    ("Global fold preserved → commitment is not global unfolding",                    False, NAVY, 1),
    ("First characterisation of the pre-commitment contact ensemble",                      True,  AMBER),
], extra=Inches(0.65))
callout(s, "This baseline is required before enhanced-sampling calculations can be designed", AMBER)
slide_num(s, 17)

# ── S18  Implications ─────────────────────────────────────────────────────────
s = S()
bar(s, "Loop CD/EF Is the Correct Engineering Target for Faster BLG Adsorption",
    AMBER, T_IMPL)
buls(s, [
    ("Loop CD/EF: the engineering target",                                                  True,  GOLD),
    ("Becomes dominant near AWI — interface-induced conformational preference",        False, NAVY, 1),
    ("Calyx flexibility mutations more effective than bulk hydrophobicity changes",         False, NAVY, 1),
    ("Dimer complexity (BLG dominant dimer above ∼50 µM):",                        True,  NAVY),
    ("Steric: dimer interface partially occludes the Loop CD/EF region",                   False, NAVY, 1),
    ("Kinetic: dimer must dissociate before calyx can present unobstructed",               False, NAVY, 1),
    ("Orientational: oblate dimer shape alters rotational diffusion to the AWI",           False, NAVY, 1),
])
slide_num(s, 18)

# ── S19  What's next ──────────────────────────────────────────────────────────
s = S()
bar(s, "Next: Enhanced Sampling Along (SASA, θ) + β-Casein Comparison",
    BLUE, T_NEXT)
buls(s, [
    ("Paper 1 — submission pipeline",                                                  True,  GOLD),
    ("Zenodo upload → DOI → JCIS submission (IF ∼9)  |  target: June 2026",  False, NAVY, 1),
    ("Enhanced sampling — the mechanistic next step",                                  True,  NAVY),
    ("Metadynamics / umbrella sampling along (SASA, θ) collective variables",         False, NAVY, 1),
    ("SET 1D vacuum-side starts as comparison baseline",                                    False, NAVY, 1),
    ("Paper 2 — β-Casein at AWI",                                                 True,  NAVY),
    ("AlphaFold2 structure ready; intrinsically disordered; adsorbs much faster",          False, NAVY, 1),
    ("Research question: why is β-casein so much faster than BLG?",                   False, BLUE, 1),
])
slide_num(s, 19)

# ── S20  Conclusions (white bg — stays on screen during Q&A) ──────────────────
s = S()
# Slim navy header (no bar() — custom layout)
rect(s, 0, 0, W, Inches(0.72), NAVY)
rect(s, 0, Inches(0.72), W, Pt(4), GOLD)
txt(s, "Conclusions", MAR, Inches(0.08), BW, Inches(0.62), size=30, color=WHITE, bold=True)

CONCL = [
    (EMERALD, "1.  Contact without commitment",
     "613 events across 4 µs — only 6 last ≥10 ns, none commits. "
     "BLG touches the AWI frequently; adsorption is not sampled by unbiased MD."),
    (OCEAN,   "2.  Compact globally, loop-mediated locally",
     "Rg stable at 1.496 ± 0.009 nm; β-barrel intact across all 4 µs. "
     "The AWI selectively amplifies Loop CD/EF mobility (0.39 nm) over Loop BC (0.54 nm in bulk)."),
    (AMBER,   "3.  SASA ⊥ orientation on the µs timescale",
     "Pearson r = +0.006, block bootstrap CI [−0.09, +0.11]. "
     "No two-factor gate. Commitment step requires enhanced sampling."),
]
y = Inches(0.86)
for color, heading, body in CONCL:
    rect(s, MAR - Inches(0.05), y, BW + Inches(0.10), Inches(1.48), LGRAY)
    rect(s, MAR - Inches(0.05), y, Pt(6), Inches(1.48), color)      # coloured left border
    txt(s, heading, MAR + Pt(12), y + Inches(0.09), BW - Pt(12), Inches(0.38),
        size=18, color=color, bold=True)
    txt(s, body, MAR + Pt(12), y + Inches(0.46), BW - Pt(12), Inches(0.92),
        size=15, color=NAVY)
    y += Inches(1.58)

txt(s, "4 µs of unbiased MD resolves the pre-commitment contact ensemble. "
       "Commitment requires enhanced sampling.",
    MAR, H - Inches(0.72), BW * 0.65, Inches(0.55), size=13, color=GRAY)
txt(s, "Questions?",
    W - MAR - Inches(3.1), H - Inches(0.92), Inches(3.0), Inches(0.50),
    size=26, color=NAVY, bold=True, align=PP_ALIGN.RIGHT)
txt(s, "github.com/cpornjar/blg-airwater-interface-md",
    W - MAR - Inches(3.1), H - Inches(0.46), Inches(3.0), Inches(0.35),
    size=12, color=BLUE, align=PP_ALIGN.RIGHT)
slide_num(s, 20)

# ── Save & deduplicate ZIP ────────────────────────────────────────────────────
prs.save(str(OUT))
seen = {}
with zipfile.ZipFile(OUT) as z:
    for item in z.infolist():
        seen[item.filename] = z.read(item.filename)
with zipfile.ZipFile(OUT, "w", zipfile.ZIP_DEFLATED) as z:
    for name, data in seen.items():
        z.writestr(name, data)
print(f"Saved: {OUT.name}  ({OUT.stat().st_size // 1024} KB)  {len(prs.slides)} slides")

# ── Push to Mac ───────────────────────────────────────────────────────────────
if shutil.which("rsync"):
    r = subprocess.run(["rsync", "-q", str(OUT), f"{MAC_HOST}:{MAC_PATH}"],
                       capture_output=True)
    if r.returncode == 0:
        subprocess.run(["ssh", MAC_HOST, f"chflags nohidden '{MAC_PATH}'"],
                       capture_output=True)
        print("Pushed to Mac and unhidden.")
