"""
make_gpu_benchmark_slide.py
============================
Generates a single-slide PowerPoint benchmarking the KU HPC GPU partitions
for the COMFHA GROMACS workload (~97k atom BLG slab).

Output: Workspace/MILK_FROTHING/GPU_BENCHMARK_SLIDE.pptx
"""

import io
from pathlib import Path
import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
import matplotlib.patches as mpatches
import numpy as np
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

ROOT  = Path(__file__).resolve().parent.parent
OUT   = ROOT / "GPU_BENCHMARK_SLIDE.pptx"

# ── Data (benchmarked partitions only) ───────────────────────────────────────
GPUS        = ["gpu1\n(NVIDIA)", "gpu2\n(NVIDIA)", "gpu-amd\n(AMD)"]
NS_PER_DAY  = [107.1, 145.9, 52.0]
RISKS       = ["LOW", "LOW–MED", "HIGH"]
COLORS      = ["#2E86AB", "#28A745", "#DC3545"]   # blue, green, red
RISK_COLORS = ["#28A745", "#FFC107", "#DC3545"]

# ── Slide dimensions (widescreen 16:9) ───────────────────────────────────────
prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

def rgb(hex_str):
    h = hex_str.lstrip("#")
    return RGBColor(int(h[0:2],16), int(h[2:4],16), int(h[4:6],16))

def add_textbox(slide, left, top, width, height,
                text, size, bold=False, color="#222222",
                align=PP_ALIGN.LEFT, wrap=True):
    tb = slide.shapes.add_textbox(Inches(left), Inches(top),
                                   Inches(width), Inches(height))
    tf = tb.text_frame
    tf.word_wrap = wrap
    p  = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = rgb(color)
    return tb

# ── Background ────────────────────────────────────────────────────────────────
bg = slide.shapes.add_shape(
    1, 0, 0, prs.slide_width, prs.slide_height)
bg.fill.solid()
bg.fill.fore_color.rgb = rgb("#F8F9FA")
bg.line.fill.background()

# ── Header bar ────────────────────────────────────────────────────────────────
hdr = slide.shapes.add_shape(
    1, 0, 0, prs.slide_width, Inches(1.05))
hdr.fill.solid()
hdr.fill.fore_color.rgb = rgb("#1B3A5C")
hdr.line.fill.background()

add_textbox(slide, 0.3, 0.08, 9, 0.55,
            "KU HPC GPU Benchmark — COMFHA Project",
            22, bold=True, color="#FFFFFF", align=PP_ALIGN.LEFT)
add_textbox(slide, 0.3, 0.58, 9, 0.35,
            "GROMACS 2020.4  |  BLG slab ~97,000 atoms  |  CHARMM36m / TIP3P  |  May 2026",
            11, bold=False, color="#AED6F1", align=PP_ALIGN.LEFT)

# ── Build bar chart with matplotlib ──────────────────────────────────────────
fig, ax = plt.subplots(figsize=(6.5, 4.2))
fig.patch.set_facecolor("#F8F9FA")
ax.set_facecolor("#F8F9FA")

x = np.arange(len(GPUS))
bars = ax.bar(x, NS_PER_DAY, color=COLORS, width=0.5,
              edgecolor="white", linewidth=1.5, zorder=3)

# Value labels
for bar, val in zip(bars, NS_PER_DAY):
    ax.text(bar.get_x() + bar.get_width()/2,
            bar.get_height() + 2.5,
            f"{val:.1f}", ha="center", va="bottom",
            fontsize=13, fontweight="bold", color="#222222")

# Reference line at gpu1 average
ax.axhline(107.1, color="#2E86AB", linewidth=1.2,
           linestyle="--", alpha=0.6, zorder=2)
ax.text(2.3, 109.5, "gpu1 baseline", color="#2E86AB",
        fontsize=8.5, ha="right")

ax.set_xticks(x)
ax.set_xticklabels(GPUS, fontsize=12, fontweight="bold")
ax.set_ylabel("ns / day", fontsize=12)
ax.set_ylim(0, 175)
ax.set_title("Simulation Speed (ns/day)", fontsize=13,
             fontweight="bold", pad=10, color="#1B3A5C")
ax.yaxis.grid(True, linestyle="--", alpha=0.5, zorder=0)
ax.set_axisbelow(True)
ax.spines["top"].set_visible(False)
ax.spines["right"].set_visible(False)
plt.tight_layout(pad=0.5)

img_buf = io.BytesIO()
fig.savefig(img_buf, format="png", dpi=150, bbox_inches="tight",
            facecolor="#F8F9FA")
plt.close(fig)
img_buf.seek(0)

slide.shapes.add_picture(img_buf, Inches(0.3), Inches(1.2),
                          Inches(6.8), Inches(4.7))

# ── Risk & notes panel (right side) ──────────────────────────────────────────
panel = slide.shapes.add_shape(
    1, Inches(7.35), Inches(1.2), Inches(5.65), Inches(5.6))
panel.fill.solid()
panel.fill.fore_color.rgb = rgb("#FFFFFF")
panel.line.color.rgb = rgb("#DEE2E6")
panel.line.width = Pt(0.75)

add_textbox(slide, 7.55, 1.30, 5.2, 0.4,
            "Risk Summary", 13, bold=True, color="#1B3A5C")

RISK_DATA = [
    ("gpu1",    "107 ns/day", "LOW",    "#28A745",
     "Proven stable · CentOS 7 · Standard GROMACS 2020.4\nRecommended default for long runs"),
    ("gpu2",    "146 ns/day", "LOW–MED","#FFC107",
     "Fastest (36% above gpu1) · 2 GPUs/node for PME offload\nSame software stack as gpu1 · More queue contention"),
    ("gpu-amd", "~52 ns/day", "HIGH",   "#DC3545",
     "2.5× slower · RHEL 9 PREEMPT kernel → job killed without\nwarning · Custom HIP binary incompatible with std. checkpoints\nR2 & R3 AMD attempts both failed; R1 died twice"),
]

y = 1.78
for name, speed, risk, rcol, note in RISK_DATA:
    # Coloured risk badge
    badge = slide.shapes.add_shape(
        1, Inches(7.55), Inches(y), Inches(1.05), Inches(0.28))
    badge.fill.solid()
    badge.fill.fore_color.rgb = rgb(rcol)
    badge.line.fill.background()
    badge_tf = badge.text_frame
    badge_tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    badge_run = badge_tf.paragraphs[0].add_run()
    badge_run.text = risk
    badge_run.font.size = Pt(8.5)
    badge_run.font.bold = True
    badge_run.font.color.rgb = rgb("#FFFFFF" if risk != "LOW–MED" else "#333333")

    # GPU name + speed
    add_textbox(slide, 8.72, y - 0.01, 4.1, 0.3,
                f"{name}  ·  {speed}", 10, bold=True, color="#222222")
    # Note
    add_textbox(slide, 7.55, y + 0.29, 5.25, 0.7,
                note, 8.5, bold=False, color="#555555")
    y += 1.12

# ── Footer ────────────────────────────────────────────────────────────────────
footer = slide.shapes.add_shape(
    1, 0, Inches(6.98), prs.slide_width, Inches(0.52))
footer.fill.solid()
footer.fill.fore_color.rgb = rgb("#1B3A5C")
footer.line.fill.background()

add_textbox(slide, 0.3, 7.02, 8, 0.38,
            "gpu3 / gpu4 excluded — no benchmark data from COMFHA workload",
            9, bold=False, color="#AED6F1", align=PP_ALIGN.LEFT)
add_textbox(slide, 9.0, 7.02, 4.1, 0.38,
            "COMFHA · KU HPC · May 2026",
            9, bold=False, color="#AED6F1", align=PP_ALIGN.RIGHT)

prs.save(str(OUT))
print(f"Saved → {OUT}")
